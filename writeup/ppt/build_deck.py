"""Build the review deck: one figure per slide, high resolution, minimal text.

The brief was a deck the professor can page through and approve -- every plot at
full resolution, one per slide, no prose. Each slide therefore carries a title
saying what the figure is, a single line saying what it shows, and the figure
itself at native resolution scaled to fit. Nothing else.

The slide order follows the manuscript's argument rather than the directory
layout: what the sensor records, that the rhythms are the ones the PSG records,
what rate accuracy that buys, the harmonic structure, the cortical negative, and
then the four channel and variance results added most recently.

Figures are embedded from `writeup/figures/` as-is, so the deck is only as
current as the last regeneration -- rerun the figure scripts first if in doubt.
Every path is checked before anything is written, so a renamed figure fails
loudly here instead of producing a deck with a slide missing.

Run from the repo root:  .venv/Scripts/python.exe writeup/ppt/build_deck.py
Output -> writeup/ppt/CAP_sleep_mask_review_deck.pptx
"""

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIGS = ROOT / "writeup" / "figures"
OUT = Path(__file__).resolve().parent / "CAP_sleep_mask_review_deck.pptx"

W_IN, H_IN = 13.333, 7.5                 # 16:9
INK = RGBColor(0x1B, 0x2A, 0x41)         # titles
MUTED = RGBColor(0x5A, 0x64, 0x72)       # one-line captions
FAINT = RGBColor(0xA8, 0xB0, 0xBA)       # provenance / slide number
ACCENT = RGBColor(0xC0, 0x39, 0x2B)      # section rules

MARGIN = 0.45
TITLE_TOP = 0.34
CAPTION_TOP = 0.86
IMG_TOP = 1.34
IMG_BOTTOM = 7.16


# ── the deck ────────────────────────────────────────────────────────────────
# (section, title, one-line caption, figure path relative to writeup/figures)
SLIDES = [
    ("What the mask records", None, None, None),
    (None, "The raw capacitive signal",
     "Both rhythms are visible in the unprocessed trace, before any filtering.",
     "signal_validation/fig1_waveform_example.png"),
    (None, "In-band signal-to-noise",
     "Respiratory and cardiac bands stand above the local noise floor on every night.",
     "signal_validation/fig2_inband_snr.png"),
    (None, "A night in the frequency domain",
     "Both rhythms persist as continuous ridges across the whole recording.",
     "spectrograms/S1N1_spectrogram_ridges.png"),

    ("Are these the rhythms the PSG records?", None, None, None),
    (None, "Coherence with the PSG sensors",
     "At the breathing and heartbeat frequencies, against shift, reverse and EEG controls.",
     "coupling/cap_psg_coherence.png"),
    (None, "Every capacitive channel against every PSG sensor",
     "Cardiac amplitude coupling is specific; respiratory amplitude largely is not.",
     "coupling/cap_psg_matrix.png"),

    ("How well the rates can be read", None, None, None),
    (None, "One night against the polysomnograph",
     "Loose peak counting on CRE, calibrated on the night being scored.",
     "rate_rerun/fig_representative_night.png"),
    (None, "Respiration — all twelve nights",
     "Reference in black, estimator after per-session calibration.",
     "rate_rerun/fig_sessions_resp.png"),
    (None, "Cardiac — all twelve nights",
     "The same presentation, cardiac band.",
     "rate_rerun/fig_sessions_card.png"),
    (None, "Agreement and its limits",
     "Bias is near zero in both bands; the limits of agreement are wide.",
     "rate_rerun/fig_bland_altman.png"),
    (None, "Accuracy by sleep stage",
     "Error is lowest in the consolidated stages and worst in wake and REM.",
     "mask_rate_detection/fig3_per_stage_mae.png"),
    (None, "Estimator against channel",
     "No channel is reliably better than the CLE−CRE differential for respiration.",
     "mask_rate_detection/fig18_mae_heatmap.png"),

    ("Harmonic structure", None, None, None),
    (None, "Persistent ridges across a night",
     "Tracked respiratory and cardiac ridges on a single channel, with the hypnogram.",
     "harmonics/ridges/ridge_tune_S1N1_CRE.png"),
    (None, "Ridge structure by sleep stage",
     "Ridge count, power and lowest frequency all separate by stage.",
     "harmonics/band_ridge_by_stage.png"),
    (None, "A harmonic comb episode",
     "Flat, richly harmonic episodes lasting tens of minutes, here in consolidated N2.",
     "harmonics/ladders/ladder_S6N1.png"),
    (None, "When the comb episodes occur",
     "They sit in N2 and tend to follow REM rather than accompany slow-wave sleep.",
     "harmonics/ladder_stage_relationship.png"),

    ("The cortical events — a mechanical response", None, None, None),
    (None, "Response at sleep spindles",
     "A clear low-band response, but no sigma-band signature: mechanical, not electrical.",
     "spindles/fig_spindle_lowband_detection.png"),
    (None, "Response at delta-burst onset",
     "The mask responds within seconds of the EEG burst, across all three channels.",
     "delta_onset/fig_delta_onset_cohort.png"),
    (None, "Capacitive against contact EEG for slow-wave activity",
     "The same pipeline that works on EEG does not transfer to the capacitive signal.",
     "swa/swa_validation_paper.png"),

    ("Channels, drift and variance", None, None, None),
    (None, "CH against the CLE−CRE differential",
     "The two differ in scale and are only partly coherent — they are not interchangeable.",
     "mean_value/ch_vs_diff_relationship.png"),
    (None, "Overnight evolution of the sensor value",
     "The DC level drifts through the night and carries a directional imbalance.",
     "channel_evolution/S1N1_CH_CLE-CRE.png"),
    (None, "Sensor value over the night — all twelve",
     "Each night referenced to its own session mean.",
     "channel_evolution/ch_vs_clecre_grid.png"),
    (None, "Capacitance imbalance — all twelve",
     "Magnitude and direction of the left-right offset across each night.",
     "imbalance/imbalance_grid.png"),
    (None, "Band amplitude by sleep stage",
     "Amplitude falls from wake into deep sleep and returns in REM, in every subject.",
     "mean_value/variance_by_stage.png"),
    (None, "How much of that variance is the instrument",
     "Against an unworn recording on the same mask: respiration clears the floor, cardiac barely.",
     "mean_value/baseline_variance_floor.png"),
    (None, "Where the high-variance epochs fall",
     "Top-decile variance per recording, with the hypnogram above and motion marked.",
     "mean_value/high_variance_traces_2col.png"),
    (None, "Which stages they fall in",
     "Observed over expected occupancy — 1.0 is chance. Almost none land in N3.",
     "mean_value/high_variance_enrichment.png"),
]

TITLE = ("Capacitive sleep mask", "Figures for review",
         "One figure per slide, at full resolution")


def textbox(slide, text, left, top, width, height, size, color,
            bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def rule(slide, left, top, width):
    """A short solid accent bar, no outline."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(0.045))
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def place_image(slide, png):
    """Scale to fit the content box, preserving aspect, centred."""
    w_px, h_px = Image.open(png).size
    box_w, box_h = W_IN - 2 * MARGIN, IMG_BOTTOM - IMG_TOP
    scale = min(box_w / w_px, box_h / h_px)
    w, h = w_px * scale, h_px * scale
    slide.shapes.add_picture(str(png), Inches((W_IN - w) / 2),
                             Inches(IMG_TOP + (box_h - h) / 2),
                             Inches(w), Inches(h))


def main():
    figs = [(t, c, f) for _, t, c, f in SLIDES if f]
    missing = [f for _, _, f in figs if not (FIGS / f).exists()]
    if missing:
        sys.exit("figures not found:\n  " + "\n  ".join(missing))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W_IN), Inches(H_IN)
    blank = prs.slide_layouts[6]

    # ── title slide
    s = prs.slides.add_slide(blank)
    textbox(s, TITLE[0], MARGIN, 2.55, W_IN - 2 * MARGIN, 1.0, 40, INK, bold=True)
    textbox(s, TITLE[1], MARGIN, 3.45, W_IN - 2 * MARGIN, 0.6, 22, MUTED)
    rule(s, MARGIN, 4.22, 1.6)
    textbox(s, TITLE[2], MARGIN, 4.45, W_IN - 2 * MARGIN, 0.5, 13, FAINT)

    n = 0
    for section, title, caption, fig in SLIDES:
        if section:                                   # ── section divider
            s = prs.slides.add_slide(blank)
            rule(s, MARGIN, 3.32, 1.1)
            textbox(s, section, MARGIN, 3.55, W_IN - 2 * MARGIN, 1.0, 30, INK, bold=True)
            continue

        s = prs.slides.add_slide(blank)
        n += 1
        textbox(s, title, MARGIN, TITLE_TOP, W_IN - 2 * MARGIN, 0.55, 23, INK, bold=True)
        textbox(s, caption, MARGIN, CAPTION_TOP, W_IN - 2 * MARGIN, 0.42, 13.5, MUTED)
        place_image(s, FIGS / fig)
        textbox(s, fig, MARGIN, 7.16, 8.0, 0.3, 8, FAINT)
        textbox(s, str(n), W_IN - MARGIN - 1.0, 7.16, 1.0, 0.3, 9, FAINT,
                align=PP_ALIGN.RIGHT)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    mb = OUT.stat().st_size / 1e6
    print("wrote %s" % OUT)
    print("  %d slides (%d figures, %d sections), %.1f MB"
          % (len(prs.slides._sldIdLst), n, sum(1 for x, *_ in SLIDES if x), mb))


if __name__ == "__main__":
    main()
